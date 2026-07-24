from __future__ import annotations
"""
Zhangl Agent - Web API
FastAPI backend with SSE streaming for agent chat.
"""

import json
import os
import re
import shutil
import sys
import time
import uuid
import base64
import io
import mimetypes
import tempfile

from fastapi import FastAPI, Request, HTTPException, Query, UploadFile, File, Form
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse, StreamingResponse, Response
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# Setup path for imports
project_root = os.path.dirname(os.path.dirname(__file__))

# Temp export directory — cleaned periodically, won't bloat disk
sys.path.insert(0, project_root)

from cli.settings import SETTINGS_DIR, load_settings, get_active_provider, get_user_dir
TEMP_EXPORT_DIR = os.path.join(SETTINGS_DIR, "tmp")
os.makedirs(TEMP_EXPORT_DIR, exist_ok=True)
from cli.repl import Repl


def _load_thinking_turns() -> int:
    try:
        return load_settings().model.expert_thinking_turns
    except Exception:
        return 2

app = FastAPI(title="Zhangl Agent", description="AI-powered agent framework")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["capacitor://localhost", "http://localhost", "http://127.0.0.1"],
    allow_credentials=False,
    allow_methods=["GET", "POST", "DELETE", "OPTIONS"],
    allow_headers=["Content-Type"],
)

# Global agent instance (initialized lazily)
_agent: Repl | None = None

# Shared AgentRegistry instance for expert sub-agents
_agent_registry = None

def _init_agent_registry():
    """Create and return an AgentRegistry, discovering built-in + custom agents."""
    global _agent_registry
    if _agent_registry is not None:
        return _agent_registry
    from agent.sub_agents.agent_registry import AgentRegistry
    _agent_registry = AgentRegistry()
    _agent_registry.auto_discover(
        os.path.join(project_root, "agent", "sub_agents"),
        os.path.join(project_root, "custom-agents"),
    )
    return _agent_registry

# Projects directory: user data lives in ~/.zhangl-agent/projects/
PROJECTS_DIR = os.path.join(get_user_dir(), "projects")


def _get_projects_dir():
    os.makedirs(PROJECTS_DIR, exist_ok=True)
    return PROJECTS_DIR


def _sanitize_name(name: str) -> str:
    return re.sub(r'[^a-zA-Z0-9_一-鿿-]', '_', name.strip())[:64]


_DEFAULT_KNOWLEDGE_TREE = {
    "判断推理": {
        "逻辑判断": ["直言命题","假言推理","选言命题","削弱加强","解释评价"],
        "图形推理": ["位置规律","样式规律","属性规律","数量规律","空间重构"],
        "定义判断": ["关键词匹配","要件提取","多定义辨析"],
        "类比推理": ["语义关系","逻辑关系","语法关系","字符类比"]
    },
    "言语理解": {
        "逻辑填空": ["语境分析","词语辨析","成语运用","关联词搭配"],
        "片段阅读": ["主旨概括","意图判断","细节理解","词句理解","标题选择"],
        "语句表达": ["语句排序","语句填空","接语选择","病句辨析"]
    },
    "资料分析": {
        "核心概念": ["增长量","增长率","比重","倍数","平均数"],
        "速算技巧": ["估算法","直除法","特殊值法","差分法","十字交叉"],
        "综合分析": ["文字材料","表格材料","图形材料","综合材料"]
    },
    "数量关系": {
        "数学运算": ["工程问题","行程问题","利润问题","排列组合","概率问题","几何问题"],
        "数字推理": ["等差数列","等比数列","幂次数列","递推数列"]
    },
    "常识判断": {
        "时政热点": ["重大会议","重要政策","领导人讲话"],
        "法律常识": ["宪法","行政法","民法典","刑法"],
        "人文历史": ["中国古代史","文学常识","传统节日"],
        "科技地理": ["前沿科技","地理国情","生活常识"]
    },
    "申论": {
        "归纳概括": ["概括主旨","提炼要点","归纳原因","总结影响"],
        "综合分析": ["词句理解","观点评价","现象分析","关系分析"],
        "提出对策": ["问题识别","对策设计","可行性分析","优先级排序"],
        "公文写作": ["通知","报告","意见","简报","倡议书"],
        "申发论述": ["立意审题","结构布局","论证方法","语言表达","素材运用"]
    }
}


def _create_default_knowledge_tree(proj_path: str):
    """Create 知识体系.json if it doesn't exist."""
    kp_path = os.path.join(proj_path, "知识体系.json")
    if not os.path.exists(kp_path):
        with open(kp_path, "w", encoding="utf-8") as f:
            json.dump(_DEFAULT_KNOWLEDGE_TREE, f, ensure_ascii=False, indent=2)


def _create_study_plan(proj_path: str, exam_date: str = "", exam_name: str = "",
                       exam_type: str = "", province: str = "", mock_count: int = 120,
                       business_model: dict | None = None):
    """Create 备考计划.json if it doesn't exist."""
    plan_path = os.path.join(proj_path, "备考计划.json")
    if not os.path.exists(plan_path):
        plan = {
            "exam_name": exam_name,
            "exam_date": exam_date,
            "exam_type": exam_type,
            "province": province,
            "mock_exam_count": mock_count,
            "business_model": business_model or {},
            "phases": {},
            "tasks": {}
        }
        if exam_date:
            # Auto-generate phases based on exam date
            from datetime import datetime, timedelta
            try:
                end = datetime.strptime(exam_date, "%Y-%m-%d")
                days = (end - datetime.now()).days
                if days > 0:
                    p1 = datetime.now().strftime("%Y-%m-%d")
                    p2 = (end - timedelta(days=days // 2)).strftime("%Y-%m-%d")
                    plan["phases"] = {"基础期": p1, "强化期": p2, "冲刺期": exam_date}
            except Exception:
                pass
        with open(plan_path, "w", encoding="utf-8") as f:
            json.dump(plan, f, ensure_ascii=False, indent=2)


class AttachmentInfo(BaseModel):
    name: str
    type: str = "file"  # "file" | "image"
    text: str = ""
    data_uri: str = ""


class ChatRequest(BaseModel):
    message: str
    session_id: str = ""
    project: str = "公考练习"  # default to exam practice project
    mode: str = "exam"  # "exam" only (QA mode archived)
    attachments: list[AttachmentInfo] = []
    thinking_mode: str = "disabled"


# ── File upload helpers ─────────────────────────────────

ALLOWED_EXTENSIONS = {".txt", ".md", ".json", ".yaml", ".yml", ".csv",
                      ".pdf", ".docx", ".xlsx", ".pptx",
                      ".png", ".jpg", ".jpeg", ".gif", ".webp"}

def _extract_file_content(filename: str, data: bytes) -> dict:
    ext = os.path.splitext(filename)[1].lower()
    mime = mimetypes.guess_type(filename)[0] or "application/octet-stream"
    result = {"type": "file", "text": "", "data_uri": None, "preview": "", "mime": mime}

    if mime.startswith("image/"):
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(data))
            max_dim = 2048
            if max(img.size) > max_dim:
                ratio = max_dim / max(img.size)
                img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.LANCZOS)
            buf = io.BytesIO()
            fmt = "PNG" if mime == "image/png" else "JPEG"
            img.save(buf, format=fmt, optimize=True)
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            result["data_uri"] = f"data:{mime};base64,{b64}"
            result["type"] = "image"
            result["preview"] = f"[Image: {filename}, {img.size[0]}x{img.size[1]}]"
        except Exception as e:
            result["preview"] = f"[Image load error: {e}]"
        return result

    if mime == "application/pdf" or ext == ".pdf":
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(data))
            lines = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    lines.append(text)
            full = "\n".join(lines)
            result["text"] = full[:50000]
            result["preview"] = f"[PDF: {filename}, {len(reader.pages)} pages, {len(data)} bytes]"
        except Exception as e:
            result["preview"] = f"[PDF parse error: {e}]"
        return result

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full = "\n".join(paragraphs)
            result["text"] = full[:50000]
            result["preview"] = f"[DOCX: {filename}, {len(doc.paragraphs)} paragraphs, {len(data)} bytes]"
        except ImportError:
            result["text"] = f"[无法解析 DOCX 文件: python-docx 未安装。请运行: pip install python-docx]"
            result["preview"] = f"[DOCX: {filename}, library not installed]"
        except Exception as e:
            result["text"] = f"[DOCX 解析失败: {e}]"
            result["preview"] = f"[DOCX parse error: {e}]"
        return result

    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), data_only=True)
            lines = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"--- Sheet: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        lines.append(row_str)
            full = "\n".join(lines)
            result["text"] = full[:50000]
            result["preview"] = f"[XLSX: {filename}, {len(wb.sheetnames)} sheets, {len(data)} bytes]"
        except ImportError:
            result["text"] = f"[无法解析 XLSX 文件: openpyxl 未安装。请运行: pip install openpyxl]"
            result["preview"] = f"[XLSX: {filename}, library not installed]"
        except Exception as e:
            result["text"] = f"[XLSX 解析失败: {e}]"
            result["preview"] = f"[XLSX parse error: {e}]"
        return result

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
            full = "\n".join(lines)
            result["text"] = full[:50000]
            result["preview"] = f"[PPTX: {filename}, {len(prs.slides)} slides, {len(data)} bytes]"
        except ImportError:
            result["text"] = f"[无法解析 PPTX 文件: python-pptx 未安装。请运行: pip install python-pptx]"
            result["preview"] = f"[PPTX: {filename}, library not installed]"
        except Exception as e:
            result["text"] = f"[PPTX 解析失败: {e}]"
            result["preview"] = f"[PPTX parse error: {e}]"
        return result

    # Text files
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = data.decode("gbk")
        except Exception:
            text = f"(binary file, {len(data)} bytes)"
    result["text"] = text[:50000]
    result["preview"] = f"[{ext.upper()}: {filename}, {len(data)} bytes]"
    return result


def _build_user_message(text: str, attachments: list[AttachmentInfo], project: str = "") -> str | list[dict]:
    has_images = any(a.type == "image" and a.data_uri for a in attachments)
    has_text_files = any(a.type == "file" and a.text for a in attachments)

    text_parts = []
    # Inject project context so AI knows where to save
    if project:
        text_parts.append(f"[参考: 用户当前打开的项目是 ~/.zhangl-agent/projects/{project}/。如果用户消息/附件是全新的需求文档，请在 ~/.zhangl-agent/projects/ 下创建新项目，不要追加到已有项目]\n")
    if has_text_files:
        text_parts.append("[用户上传了以下文件，内容已提取，直接阅读下面的内容即可，无需再 read_file]\n")
        for a in attachments:
            if a.type == "file" and a.text:
                text_parts.append(f"=== 文件: {a.name} ===\n{a.text}\n")
        text_parts.append("=== 用户消息 ===")
    text_parts.append(text)
    full_text = "\n".join(text_parts)

    if not has_images:
        return full_text

    content_blocks = [{"type": "text", "text": full_text}]
    for a in attachments:
        if a.type == "image" and a.data_uri:
            content_blocks.append({
                "type": "image_url",
                "image_url": {"url": a.data_uri}
            })
    return content_blocks


class NewSessionRequest(BaseModel):
    name: str = ""


@app.post("/api/chat/stop")
async def stop_chat():
    """Stop the currently running agent generation."""
    agent = get_agent()
    if agent.engine:
        agent.engine.cancel()
    return {"message": "Stop signal sent"}


class SettingsSaveRequest(BaseModel):
    provider: str = "deepseek"
    api_key: str = ""
    api_base: str = ""
    model: str = ""
    max_tokens: int = 32768
    fast_model: str = ""
    smart_model: str = ""
    expert_thinking_turns: int = 2


@app.get("/api/settings/status")
async def settings_status():
    """Check if the app has been configured with an API key."""
    from cli.settings import load_settings as _load, get_active_provider
    try:
        settings = _load()
        provider_type, api_key, api_base, model, max_tokens = get_active_provider(settings)
        configured = bool(api_key)
        mc = settings.model
        return {
            "configured": configured,
            "provider": mc.model_provider if configured else "",
            "model": {
                "MODEL_PROVIDER": mc.model_provider,
                "ZHANGL_BASE_URL": mc.base_url,
                "ZHANGL_AUTH_TOKEN": mc.auth_token if configured else "",
                "DEFAULT_MODEL": mc.default_model,
                "DEFAULT_MODEL_MAX_TOKENS": mc.default_model_max_tokens,
                "SMART_MODEL": mc.smart_model,
                "SMART_MODEL_MAX_TOKENS": mc.smart_model_max_tokens,
                "SMALL_MODEL": mc.small_model,
                "SMALL_MODEL_MAX_TOKENS": mc.small_model_max_tokens,
                "EXPERT_THINKING_TURNS": mc.expert_thinking_turns,
            } if configured else {},
            "model_name": model if configured else "",
        }
    except Exception:
        return {"configured": False, "provider": "", "model": {}}


@app.post("/api/settings/save")
async def settings_save(req: SettingsSaveRequest):
    """Save provider settings. Direct file write for cross-platform reliability."""
    try:
        import json, sys, os, traceback

        if sys.platform == "win32":
            settings_dir = os.path.join(os.environ.get("APPDATA", os.path.expanduser("~")), "zhangl-agent")
        else:
            settings_dir = get_user_dir()
        settings_file = os.path.join(settings_dir, "settings.json")

        print(f"[settings_save] platform={sys.platform}, dir={settings_dir}", flush=True)

        os.makedirs(settings_dir, exist_ok=True)

        data = {
            "model": {
                "MODEL_PROVIDER": req.provider,
                "ZHANGL_BASE_URL": req.api_base,
                "ZHANGL_AUTH_TOKEN": req.api_key,
                "DEFAULT_MODEL": req.model or "deepseek-v4-flash",
                "DEFAULT_MODEL_MAX_TOKENS": req.max_tokens or 32768,
                "SMART_MODEL": req.smart_model or "",
                "SMART_MODEL_MAX_TOKENS": 32768,
                "SMALL_MODEL": req.fast_model or "",
                "SMALL_MODEL_MAX_TOKENS": 32768,
                "EXPERT_THINKING_TURNS": req.expert_thinking_turns,
            },
            "permissions": {"allow": [], "deny": [], "ask_before": ["write_file"]},
            "memory": {"enabled": True, "storage_dir": "", "auto_remember": True},
            "export": {"default_format": "json", "default_dir": "./test_cases",
                       "formats": {"json": {"enabled": True}, "excel": {"enabled": True},
                                   "markdown": {"enabled": False}, "testrail_csv": {"enabled": False}}},
            "status_line": {"type": "default", "items": ["model", "provider", "tokens"]},
            "ui": {"theme": "dark", "dialog_style": "panel"},
        }

        with open(settings_file, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)

        print(f"[settings_save] saved OK to {settings_file}", flush=True)

        global _agent
        _agent = None
        os.environ["OPENAI_API_KEY"] = req.api_key
        if req.api_base:
            os.environ["OPENAI_API_BASE"] = req.api_base

        return {"message": "Settings saved", "provider": req.provider}
    except Exception as e:
        tb = traceback.format_exc()
        print(f"[settings_save] ERROR: {tb}", flush=True)
        raise HTTPException(status_code=500, detail=f"{type(e).__name__}: {e}\n{tb}")


class CasesSaveRequest(BaseModel):
    cases: list[dict] = []


class CaseUpdateRequest(BaseModel):
    case: dict = {}


def get_agent() -> Repl:
    global _agent
    if _agent is None:
        settings = load_settings()
        provider_type, api_key, api_base, model, max_tokens = get_active_provider(settings)
        _agent = Repl(
            settings=settings,
            api_key=api_key,
            api_base=api_base,
            model=model,
            provider=provider_type,
            max_tokens=max_tokens,
        )
        # Inject LLM/tool_registry into spawn_expert so background experts work
        from tools.core.spawn_expert import set_expert_dependencies
        from agent.sub_agents.agent_registry import AgentRegistry
        agent_registry = _init_agent_registry()
        set_expert_dependencies(_agent.llm, _agent.skill_registry, settings, agent_registry)
        # Ensure a session exists for the web backend
        if not _agent.sessions.current_id:
            _agent.sessions.create(model, provider_type)

    # Always sync model with current settings — settings may have changed since startup
    settings = load_settings()
    provider_type, api_key, api_base, model, max_tokens = get_active_provider(settings)
    if _agent.model != model or _agent.api_key != api_key:
        _agent.model = model
        _agent.max_tokens = max_tokens
        _agent.api_key = api_key
        _agent.api_base = api_base
        _agent.provider = provider_type
        if hasattr(_agent, 'engine') and _agent.engine:
            _agent.engine.shutdown()
        _agent._init_agent()
        from tools.core.spawn_expert import set_expert_dependencies
        from agent.sub_agents.agent_registry import AgentRegistry
        agent_registry = _init_agent_registry()
        set_expert_dependencies(_agent.llm, _agent.skill_registry, settings, agent_registry)

    return _agent


@app.get("/")
async def index():
    static_dir = os.path.join(os.path.dirname(__file__), "static")
    with open(os.path.join(static_dir, "index.html"), encoding="utf-8") as f:
        return HTMLResponse(f.read(), headers={"Cache-Control": "no-cache, no-store, must-revalidate"})


@app.post("/api/chat/upload")
async def chat_upload(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(400, "No filename")
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported file type: {ext}. Supported: {', '.join(sorted(ALLOWED_EXTENSIONS))}")
    data = await file.read()
    if len(data) > 20 * 1024 * 1024:
        raise HTTPException(400, "File too large (max 20 MB)")
    info = _extract_file_content(file.filename, data)
    return {
        "filename": file.filename,
        "size": len(data),
        "type": info["type"],
        "text": info["text"],
        "data_uri": info["data_uri"],
        "preview": info["preview"],
    }


@app.post("/api/chat/stream")
async def chat_stream(req: ChatRequest):
    """Stream agent output as SSE events — text, tool calls, results in real time."""
    from agent.engine import TextDelta, ToolCallStart, ToolCallResult, AgentDone, AgentError
    agent = get_agent()
    agent.llm._thinking_mode = req.thinking_mode

    # Switch to the requested session if different (for fresh-session tasks)
    original_session_id = agent.sessions.current_id
    switched = False
    if req.session_id and req.session_id != original_session_id:
        try:
            agent._save_session()
            meta, msgs, summary, count = agent.sessions.resume(req.session_id)
            if meta:
                agent.model = meta.model
                agent._init_agent()
                agent.ctx_mgr.messages = msgs
                agent.ctx_mgr._summary = summary
                agent.ctx_mgr._summarized_count = count
                switched = True
        except Exception:
            pass  # stay on current session

    async def generate():
        files_created = []
        final_text = []
        user_content = _build_user_message(req.message, req.attachments, req.project)
        async for event in agent.engine.run(user_content):
            if isinstance(event, TextDelta):
                final_text.append(event.content)
                yield f"data: {json.dumps({'type': 'text', 'content': event.content}, ensure_ascii=False)}\n\n"
            elif isinstance(event, ToolCallStart):
                # Look up tool label from registry for frontend display
                label = event.name  # fallback
                try:
                    t = agent.skill_registry.get(event.name)
                    if t:
                        label = t.description.split('。')[0].split('.')[0].strip()
                        if len(label) > 50:
                            label = label[:50] + '...'
                except Exception:
                    pass
                payload = {'type': 'tool_start', 'name': event.name, 'args': event.arguments, 'label': label}
                if event.expert_id:
                    payload['expert_id'] = event.expert_id
                yield f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"
            elif isinstance(event, ToolCallResult):
                result_payload = {'type': 'tool_result', 'name': event.name, 'result': event.result[:200]}
                if event.expert_id:
                    result_payload['expert_id'] = event.expert_id
                if event.name == "spawn_expert":
                    import re
                    tools_used = re.findall(r'\[(\w+)\]:', event.result)
                    if tools_used:
                        result_payload['expert_tools'] = list(dict.fromkeys(tools_used))
                        result_payload['result'] = event.result[:500]
                if event.name in ("write_file", "export_json", "export_excel",
                                  "export_markdown", "export_xmind", "export_testrail_csv"):
                    result = event.result
                    if "File saved:" in result or "Exported" in result:
                        for line in result.split("\n"):
                            if " to " in line:
                                files_created.append(line.split(" to ")[-1].strip())
                            elif "File saved:" in line:
                                parts = line.split("(")[0].replace("File saved:", "").strip()
                                if parts:
                                    files_created.append(parts)
                elif event.name == "run_bash":
                    import re
                    cmd = event.arguments.get("command", "")
                    for pat in [r'(?:>>?|cp\s+.*?)\s*(projects/[^\s;|&]+)', r'mkdir\s+-p\s*(projects/[^\s;|&]+)']:
                        for m in re.finditer(pat, cmd):
                            fpath = m.group(1).rstrip('/')
                            if fpath not in files_created:
                                files_created.append(fpath)
                yield f"data: {json.dumps(result_payload, ensure_ascii=False)}\n\n"
            elif isinstance(event, AgentDone):
                agent._save_session()
                if switched and original_session_id:
                    summary_text = ''.join(final_text).strip()
                    if summary_text:
                        try:
                            meta2, msgs2, summ2, cnt2 = agent.sessions.resume(original_session_id)
                            if meta2:
                                agent.model = meta2.model
                                agent._init_agent()
                                agent.ctx_mgr.messages = msgs2
                                agent.ctx_mgr._summary = summ2
                                agent.ctx_mgr._summarized_count = cnt2
                                # Brief summary: first sentence, capped at 500 chars as safety
                                brief = summary_text.split('。')[0].split('\n')[0][:500]
                                agent.ctx_mgr.add_message({"role": "user", "content": f"[系统通知] 子任务已完成。{brief}"})
                                agent._save_session()
                        except Exception:
                            pass
                yield f"data: {json.dumps({'type': 'done', 'files': files_created}, ensure_ascii=False)}\n\n"
                return
            elif isinstance(event, AgentError):
                agent._save_session()
                if switched and original_session_id:
                    try:
                        meta2, msgs2, summ2, cnt2 = agent.sessions.resume(original_session_id)
                        if meta2:
                            agent.model = meta2.model
                            agent._init_agent()
                            agent.ctx_mgr.messages = msgs2
                            agent.ctx_mgr._summary = summ2
                            agent.ctx_mgr._summarized_count = cnt2
                            agent.ctx_mgr.add_message({"role": "user", "content": f"[系统通知] 子任务失败：{str(event.message)[:300]}"})
                            agent._save_session()
                    except Exception:
                        pass
                yield f"data: {json.dumps({'type': 'error', 'message': event.message}, ensure_ascii=False)}\n\n"
                return

    return StreamingResponse(generate(), media_type="text/event-stream; charset=utf-8")


@app.post("/api/chat")
async def chat(req: ChatRequest):
    """Run agent to completion and return results as JSON. No streaming."""
    agent = get_agent()
    agent.llm._thinking_mode = req.thinking_mode

    # Session switching (same as streaming endpoint)
    original_session_id = agent.sessions.current_id
    switched = False
    if req.session_id and req.session_id != original_session_id:
        try:
            agent._save_session()
            meta, msgs, summary, count = agent.sessions.resume(req.session_id)
            if meta:
                agent.model = meta.model
                agent._init_agent()
                agent.ctx_mgr.messages = msgs
                agent.ctx_mgr._summary = summary
                switched = True
        except Exception:
            pass

    summary_parts = []
    files_created = []
    error_msg = None

    try:
        user_content = _build_user_message(req.message, req.attachments, req.project)
        async for event in agent.engine.run(user_content):
            evt_name = event.__class__.__name__
            if evt_name == "TextDelta":
                summary_parts.append(event.content)
            elif evt_name == "ToolCallResult":
                if event.name in ("write_file", "export_json", "export_excel",
                                  "export_markdown", "export_xmind", "export_testrail_csv"):
                    result = event.result
                    if "File saved:" in result or "Exported" in result:
                        for line in result.split("\n"):
                            if " to " in line:
                                files_created.append(line.split(" to ")[-1].strip())
                            elif "File saved:" in line:
                                parts = line.split("(")[0].replace("File saved:", "").strip()
                                if parts:
                                    files_created.append(parts)
                elif event.name == "run_bash":
                    import re
                    cmd = event.arguments.get("command", "")
                    for pat in [r'(?:>>?|cp\s+.*?)\s*(projects/[^\s;|&]+)', r'mkdir\s+-p\s*(projects/[^\s;|&]+)']:
                        for m in re.finditer(pat, cmd):
                            fpath = m.group(1).rstrip('/')
                            if fpath not in files_created:
                                files_created.append(fpath)
            elif evt_name == "AgentError":
                error_msg = event.message
                break
    except Exception as e:
        error_msg = str(e)

    # Save session after completion, switch back if needed
    agent._save_session()
    if switched:
        try:
            agent.sessions.resume(original_session_id)
            agent._init_agent()
        except Exception:
            pass

    full_text = "".join(summary_parts).strip()
    # Show text directly, cap at reasonable display length
    if len(full_text) <= 500:
        summary = full_text
    else:
        # Find last non-JSON line for summary, or truncate
        lines = full_text.split("\n")
        for line in reversed(lines):
            line = line.strip()
            if line and not line.startswith("[") and not line.startswith("{"):
                summary = line[:300]
                break
        else:
            summary = full_text[:300]

    # Deduplicate files
    seen = set()
    unique_files = []
    for f in files_created:
        if f not in seen:
            seen.add(f)
            unique_files.append(f)

    return {
        "summary": summary or "任务完成",
        "files": unique_files,
        "error": error_msg,
    }


@app.get("/api/sessions")
async def list_sessions():
    agent = get_agent()
    sessions = agent.sessions.list_sessions(50)
    return [
        {
            "id": s.id,
            "name": s.name[:60],
            "model": s.model,
            "message_count": s.message_count,
            "updated_at": s.updated_at[:16],
        }
        for s in sessions
    ]


@app.delete("/api/sessions/{session_id}")
async def delete_session(session_id: str):
    agent = get_agent()
    # Can't delete current session — switch to a different one first
    if agent.sessions.current_id == session_id:
        raise HTTPException(400, "Cannot delete current active session")
    agent.sessions.delete(session_id)
    return {"ok": True}


@app.post("/api/sessions/clear-all")
async def clear_all_sessions():
    """Delete all non-current sessions."""
    agent = get_agent()
    current_id = agent.sessions.current_id
    sessions = agent.sessions.list_sessions(500)
    count = 0
    for s in sessions:
        if s.id != current_id:
            agent.sessions.delete(s.id)
            count += 1
    return {"ok": True, "deleted": count}


@app.post("/api/sessions/clear-all-including-current")
async def clear_all_sessions_including_current():
    """Delete ALL sessions including the current one. Creates a fresh session after."""
    agent = get_agent()
    current_id = agent.sessions.current_id
    sessions = agent.sessions.list_sessions(500)
    count = 0
    for s in sessions:
        agent.sessions.delete(s.id)
        count += 1
    if current_id:
        agent.sessions.delete(current_id)
        count += 1
    # Create a fresh session
    settings = load_settings()
    provider_type, api_key, api_base, model, max_tokens = get_active_provider(settings)
    agent.model = model
    agent.max_tokens = max_tokens
    agent._init_agent()
    meta = agent.sessions.create(agent.model, agent.provider_label.lower())
    agent.ctx_mgr = __import__("context.manager", fromlist=["ContextManager"]).ContextManager(keep_thinking_turns=_load_thinking_turns())
    agent._init_agent()
    return {"ok": True, "deleted": count, "new_session_id": meta.id}


@app.post("/api/sessions/new")
async def new_session(req: NewSessionRequest):
    agent = get_agent()
    agent._save_session()
    # Sync model with current settings — settings may have changed since startup
    settings = load_settings()
    provider_type, api_key, api_base, model, max_tokens = get_active_provider(settings)
    agent.model = model
    agent.max_tokens = max_tokens
    agent._init_agent()
    # Re-inject dependencies in case settings changed
    from tools.core.spawn_expert import set_expert_dependencies
    from agent.sub_agents.agent_registry import AgentRegistry
    agent_registry = _init_agent_registry()
    set_expert_dependencies(agent.llm, agent.skill_registry, settings, agent_registry)
    from tools.core.task_tools import reset_tasks
    reset_tasks()
    meta = agent.sessions.create(agent.model, agent.provider_label.lower(), name=req.name)
    agent.ctx_mgr = __import__("context.manager", fromlist=["ContextManager"]).ContextManager(keep_thinking_turns=_load_thinking_turns())
    agent._init_agent()
    return {"id": meta.id, "name": meta.name}


@app.post("/api/sessions/{session_id}/resume")
async def resume_session(session_id: str):
    agent = get_agent()
    agent._save_session()
    meta, msgs, summary, count = agent.sessions.resume(session_id)
    if not meta:
        raise HTTPException(404, "Session not found")
    agent.model = meta.model
    agent._init_agent()
    agent.ctx_mgr.messages = msgs
    agent.ctx_mgr._summary = summary
    agent.ctx_mgr._summarized_count = count
    agent.ctx_mgr._initialized = True
    return {"id": meta.id, "name": meta.name, "message_count": meta.message_count}


@app.get("/api/sessions/current")
async def current_session():
    """Get current session info and messages for frontend restore on page load."""
    agent = get_agent()
    if not agent.sessions.current_id:
        return {"messages": [], "session_id": ""}
    meta, msgs, summary, count = agent.sessions.load(agent.sessions.current_id)
    # Filter to user/assistant messages only (skip system/tool messages for display)
    display_msgs = []
    for m in msgs:
        if m.get("role") in ("user", "assistant") and m.get("content"):
            ct = m["content"]
            # If content is a list of blocks (thinking + text), extract text only
            if isinstance(ct, list):
                text_parts = [b.get("text", "") for b in ct if b.get("type") == "text"]
                display_text = "\n".join(text_parts)
            else:
                display_text = str(ct)
            if display_text:
                display_msgs.append({"role": m["role"], "content": display_text})
    return {
        "session_id": agent.sessions.current_id,
        "name": meta.name if meta else "",
        "messages": display_msgs,
    }


@app.get("/api/tools")
async def list_tools():
    agent = get_agent()
    return [
        {"name": t.name, "description": t.description[:100], "parameters": list(t.parameters.get("properties", {}).keys())}
        for t in agent.skill_registry.list_all()
    ]


@app.get("/api/expert-status")
async def expert_status():
    """Diagnostic endpoint: show expert subsystem state.
    Returns init source, LLM info, tool registry, expert type readiness, running tasks.
    """
    # Ensure agent is initialized (injects dependencies)
    get_agent()
    from tools.core.spawn_expert import validate_expert_system
    return validate_expert_system()


@app.post("/api/projects/{name}/export/{format}")
async def export_project_cases(name: str, format: str):
    """Direct export - no AI agent needed. Reads project test cases and exports to temp dir."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")

    cases_path = os.path.join(proj_path, "cases", "test_cases.jsonl")
    if not os.path.exists(cases_path):
        raise HTTPException(404, "No test cases found. Generate test cases first.")

    # Clean up exports older than 1 hour
    _cleanup_temp_exports()

    import uuid as _uuid
    uid = _uuid.uuid4().hex[:8]
    os.makedirs(TEMP_EXPORT_DIR, exist_ok=True)

    try:
        if format == "json":
            fname = f"{uid}_test_cases.json"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export import export_json
            result = export_json(cases_path, output)
        elif format == "excel":
            fname = f"{uid}_test_cases.xlsx"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export import export_excel
            result = export_excel(cases_path, output)
        elif format == "markdown":
            fname = f"{uid}_test_cases.md"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export import export_markdown
            result = export_markdown(cases_path, output)
        elif format == "pdf":
            fname = f"{uid}_test_cases.pdf"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export import export_pdf
            result = export_pdf(cases_path, output)
        elif format == "csv":
            fname = f"{uid}_test_cases.csv"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export import export_csv
            result = export_csv(cases_path, output)
        elif format == "testrail_csv":
            fname = f"{uid}_test_cases.csv"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export import export_testrail_csv
            result = export_testrail_csv(cases_path, output)
        elif format == "xmind":
            fname = f"{uid}_test_cases.xmind"
            output = os.path.join(TEMP_EXPORT_DIR, fname)
            from tools.export.xmind_exporter import export_xmind
            result = export_xmind(json.dumps(_load_project_cases(name), ensure_ascii=False), output)
        else:
            raise HTTPException(400, f"Unknown format: {format}")

        download_url = f"/api/exports/{fname}"
        return {"success": True, "result": result, "download_url": download_url, "filename": fname}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Export failed: {e}")


@app.get("/api/exports/{filename}")
async def download_project_export(filename: str):
    """Download an exported file from the temp export directory."""
    export_path = os.path.join(TEMP_EXPORT_DIR, filename)
    if not os.path.exists(export_path):
        raise HTTPException(404, "File not found or expired. Please re-export.")
    return FileResponse(export_path, filename=filename)


def _cleanup_temp_exports():
    """Remove exports older than 1 hour to prevent disk bloat."""
    try:
        now = time.time()
        for f in os.listdir(TEMP_EXPORT_DIR):
            fpath = os.path.join(TEMP_EXPORT_DIR, f)
            if os.path.isfile(fpath) and (now - os.path.getmtime(fpath)) > 3600:
                os.remove(fpath)
    except Exception:
        pass


# --- Project Management APIs ---

@app.get("/api/projects")
async def list_projects():
    """List all projects with stats."""
    projects_dir = _get_projects_dir()
    projects = []
    for name in sorted(os.listdir(projects_dir)):
        proj_path = os.path.join(projects_dir, name)
        if not os.path.isdir(proj_path):
            continue
        cases = _load_project_cases(name)
        passed = sum(1 for c in cases if c.get("status") == "passed")
        total = len(cases)
        pass_rate = round(passed / total * 100) if total > 0 else 0
        projects.append({
            "name": name,
            "case_count": len(cases),
            "modules": len(set(c.get("module", "") for c in cases)),
            "updated_at": time.strftime("%Y-%m-%d %H:%M", time.localtime(os.path.getmtime(proj_path))),
            "has_spec": os.path.exists(os.path.join(proj_path, "spec")),
            "pass_rate": pass_rate,
            "passed_count": passed,
        })
    return sorted(projects, key=lambda p: p["updated_at"], reverse=True)


@app.post("/api/projects")
async def create_project(request: Request, name: str = Form(""), requirements: str = Form("")):
    """Create a new project."""
    # Read body once (JSON or Form)
    json_body = {}
    if not name:
        try:
            json_body = await request.json()
            name = json_body.get("name", "") or json_body.get("exam_name", "")
            requirements = json_body.get("requirements", "")
        except Exception:
            pass
    if not name:
        name = f"project_{time.strftime('%Y%m%d_%H%M%S')}"
    name = _sanitize_name(name)
    proj_path = os.path.join(_get_projects_dir(), name)
    if os.path.exists(proj_path):
        raise HTTPException(400, "Project already exists")
    os.makedirs(os.path.join(proj_path, "spec"))
    os.makedirs(os.path.join(proj_path, "cases"))
    os.makedirs(os.path.join(proj_path, "reports"))
    if requirements:
        with open(os.path.join(proj_path, "spec", "requirements.md"), "w", encoding="utf-8") as f:
            f.write(requirements)
    # Create default knowledge tree and study plan
    _create_default_knowledge_tree(proj_path)
    _create_study_plan(proj_path, json_body.get("exam_date", ""), json_body.get("exam_name", name),
                       json_body.get("exam_type", ""), json_body.get("province", ""),
                       json_body.get("mock_exam_count", 120), json_body.get("business_model") or {})
    return {"name": name, "message": f"Project '{name}' created"}


@app.delete("/api/projects/{name}")
async def delete_project(name: str):
    """Delete a project and all its data."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")
    shutil.rmtree(proj_path)
    return {"message": f"Project '{name}' deleted"}


@app.post("/api/projects/{name}/upload")
async def upload_file(name: str, file: UploadFile = File(...)):
    """Upload a spec file to a project."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")
    spec_dir = os.path.join(proj_path, "spec")
    os.makedirs(spec_dir, exist_ok=True)
    content = await file.read()
    path = os.path.join(spec_dir, file.filename or "upload")
    with open(path, "wb") as f:
        f.write(content)
    # Read content for display
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = f"(binary file, {len(content)} bytes)"
    return {"filename": file.filename, "size": len(content), "preview": text[:3000]}


@app.get("/api/projects/{name}")
async def get_project(name: str):
    """Get project info: cases, specs, stats."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")

    cases = _load_project_cases(name)
    specs = _list_spec_files(name)

    # Compute stats
    modules = {}
    priorities = {"P0": 0, "P1": 0, "P2": 0, "P3": 0}
    for tc in cases:
        mod = tc.get("module", "Other")
        modules[mod] = modules.get(mod, 0) + 1
        p = tc.get("priority", "P2")
        priorities[p] = priorities.get(p, 0) + 1

    return {
        "name": name,
        "cases": cases,
        "case_count": len(cases),
        "modules": modules,
        "priorities": priorities,
        "spec_files": specs,
        "has_spec": len(specs) > 0,
    }


# ── Coverage analysis ───────────────────────────────────

@app.post("/api/projects/{name}/analyze-coverage")
async def analyze_project_coverage(name: str, deep: bool = Query(False)):
    """Run requirement coverage analysis for a project.

    By default runs the deterministic analyze_requirement_coverage tool directly.
    Set deep=true to trigger AI agent for intelligent gap analysis and quality assessment.
    """
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")

    if not deep:
        from tools.analyze.coverage_analyzer import analyze_requirement_coverage as _cov_func
        result_json = _cov_func(project_name=name, output_format="full")
        result = json.loads(result_json)
        return result

    # Deep mode: use AI agent for intelligent analysis
    agent = get_agent()
    prompt = (
        f"分析 projects/{name}/ 项目的测试覆盖率。\n"
        f"1. 读取 spec/ 目录下的需求文档\n"
        f"2. 读取 cases/test_cases.jsonl 中的测试用例\n"
        f"3. 分析需求到用例的追溯关系\n"
        f"4. 识别未覆盖的需求\n"
        f"5. 评估现有用例的质量\n"
        f"6. 给出具体改进建议\n"
        f"最后用中文输出一份结构化的覆盖率分析报告。"
    )
    summary_parts = []
    try:
        async for event in agent.engine.run(prompt):
            from agent.engine import TextDelta, AgentError, AgentDone
            if isinstance(event, TextDelta):
                summary_parts.append(event.content)
            elif isinstance(event, AgentError):
                return {"error": event.message, "analysis_text": "".join(summary_parts)}
            elif isinstance(event, AgentDone):
                break
    except Exception as e:
        return {"error": str(e), "analysis_text": "".join(summary_parts)}

    return {"deep_analysis": True, "analysis_text": "".join(summary_parts)}


@app.get("/api/projects/{name}/coverage-report")
async def get_coverage_report(name: str):
    """Get cached coverage report for a project."""
    report_path = os.path.join(_get_projects_dir(), name, "reports", "coverage.json")
    if not os.path.exists(report_path):
        return {"cached": False, "message": "No cached report. Run analyze-coverage first."}
    try:
        with open(report_path, encoding="utf-8") as f:
            return {"cached": True, **json.load(f)}
    except (IOError, json.JSONDecodeError):
        return {"cached": False, "message": "Report file corrupted. Re-run analyze-coverage."}


# ── Test execution ──────────────────────────────────────

@app.post("/api/projects/{name}/run-tests")
async def run_project_tests(name: str):
    """Auto-detect test framework and run tests for a project.
    If no test code found but test_cases.jsonl exists, return case metadata
    for AI-assisted execution."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")

    from tools.execute.auto_runner import detect_and_run_tests as _run_func
    result_json = _run_func(project_dir=proj_path, project_name=name)
    result = json.loads(result_json)

    # No framework detected — check if project has JSONL test cases
    if result.get("status") == "no_framework":
        cases_dir = os.path.join(proj_path, "cases")
        jsonl_files = []
        total_cases = 0
        if os.path.isdir(cases_dir):
            for f in os.listdir(cases_dir):
                if f.endswith(".jsonl") or f.endswith(".json"):
                    fp = os.path.join(cases_dir, f)
                    count = _count_cases(fp)
                    jsonl_files.append({"file": f, "cases": count})
                    total_cases += count

        if total_cases > 0:
            result["status"] = "ai_execute"
            result["reason"] = f"项目包含 {total_cases} 个用例（描述文件），需要 AI 辅助执行"
            result["case_files"] = jsonl_files
            result["total_cases"] = total_cases

    return result


@app.post("/api/projects/{name}/cases")
async def save_cases(name: str, req: CasesSaveRequest):
    """Batch save/overwrite test cases for a project."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.exists(proj_path):
        raise HTTPException(404, "Project not found")
    cases_dir = os.path.join(proj_path, "cases")
    os.makedirs(cases_dir, exist_ok=True)
    path = os.path.join(cases_dir, "test_cases.jsonl")
    cases = req.cases
    for i, tc in enumerate(cases):
        if not tc.get("id"):
            mod = tc.get("module", "MOD")
            tc["id"] = f"TC-{mod}-{i+1:03d}"
    _save_project_cases(path, cases)
    return {"saved": len(cases), "path": path}


@app.post("/api/projects/{name}/cases/update")
async def update_case(name: str, req: CaseUpdateRequest):
    """Update a single test case."""
    cases = _load_project_cases(name)
    case = req.case
    cid = case.get("id", "")
    updated = False
    for i, existing in enumerate(cases):
        if existing.get("id") == cid:
            cases[i] = case
            updated = True
            break
    if not updated:
        cases.append(case)
    proj_path = os.path.join(_get_projects_dir(), name)
    cases_dir = os.path.join(proj_path, "cases")
    os.makedirs(cases_dir, exist_ok=True)
    _save_project_cases(os.path.join(cases_dir, "test_cases.jsonl"), cases)
    return {"message": "updated" if updated else "added", "id": cid}


@app.delete("/api/projects/{name}/cases/{case_id}")
async def delete_case(name: str, case_id: str):
    """Delete a single test case."""
    cases = _load_project_cases(name)
    cases = [c for c in cases if c.get("id") != case_id]
    proj_path = os.path.join(_get_projects_dir(), name)
    cases_dir = os.path.join(proj_path, "cases")
    os.makedirs(cases_dir, exist_ok=True)
    _save_project_cases(os.path.join(cases_dir, "test_cases.jsonl"), cases)
    return {"message": f"Case {case_id} deleted"}


@app.get("/api/projects/{name}/files/{filename:path}")
async def get_project_file(name: str, filename: str):
    """Get a project file or list directory."""
    proj_path = os.path.join(_get_projects_dir(), name)
    path = os.path.join(proj_path, filename)
    if not os.path.exists(path):
        raise HTTPException(404, "File not found")
    # Directory listing
    if os.path.isdir(path):
        entries = []
        for entry in sorted(os.listdir(path)):
            entry_path = os.path.join(path, entry)
            entries.append({
                "name": entry,
                "is_dir": os.path.isdir(entry_path),
                "size": os.path.getsize(entry_path) if os.path.isfile(entry_path) else 0,
                "modified": os.path.getmtime(entry_path)
            })
        return entries
    return FileResponse(path)


@app.post("/api/projects/{name}/files/{filename:path}")
async def write_project_file(name: str, filename: str, body: dict):
    """Write content to a project file. Body: {"content": "..."}"""
    proj_path = os.path.join(_get_projects_dir(), name)
    path = os.path.join(proj_path, filename)
    # Security: prevent path traversal
    if not os.path.realpath(path).startswith(os.path.realpath(proj_path)):
        raise HTTPException(403, "Path traversal denied")
    # Only allow writes to existing projects
    if not os.path.isdir(proj_path):
        raise HTTPException(404, "Project not found")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    content = body.get("content", "")
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    return {"message": "ok", "path": filename, "size": len(content)}


@app.delete("/api/projects/{name}/files/{filename:path}")
async def delete_project_file(name: str, filename: str):
    """Delete a project file."""
    proj_path = os.path.join(_get_projects_dir(), name)
    path = os.path.join(proj_path, filename)
    if not os.path.realpath(path).startswith(os.path.realpath(proj_path)):
        raise HTTPException(403, "Path traversal denied")
    if not os.path.isfile(path):
        raise HTTPException(404, "File not found")
    os.remove(path)
    return {"message": "ok", "path": filename}


@app.get("/api/projects/{name}/export")
async def export_project(name: str):
    """Export entire project as a zip file."""
    import zipfile, io
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.isdir(proj_path):
        raise HTTPException(404, "Project not found")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(proj_path):
            for f in files:
                fpath = os.path.join(root, f)
                arcname = os.path.relpath(fpath, proj_path)
                zf.write(fpath, arcname)
    buf.seek(0)
    from urllib.parse import quote
    safe_name = quote(name) + '_export.zip'
    return Response(
        content=buf.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{safe_name}"}
    )


@app.post("/api/projects/{name}/import")
async def import_project(name: str, file: UploadFile = File(...)):
    """Import a zip file into an existing project (merge, don't overwrite)."""
    import zipfile, io
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.isdir(proj_path):
        raise HTTPException(404, "Project not found")
    content = await file.read()
    try:
        buf = io.BytesIO(content)
        with zipfile.ZipFile(buf, 'r') as zf:
            count = 0
            for info in zf.infolist():
                if info.is_dir():
                    continue
                dest = os.path.join(proj_path, info.filename)
                os.makedirs(os.path.dirname(dest), exist_ok=True)
                # Only write if file doesn't exist (merge, don't overwrite)
                if not os.path.exists(dest):
                    with zf.open(info) as src, open(dest, 'wb') as dst:
                        dst.write(src.read())
                    count += 1
        return {"message": f"导入完成，新增 {count} 个文件"}
    except zipfile.BadZipFile:
        raise HTTPException(400, "无效的 ZIP 文件")


@app.get("/api/projects/{name}/files/")
async def list_project_files(name: str):
    """List all files in a project with total size."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.isdir(proj_path):
        raise HTTPException(404, "Project not found")
    files = []
    for root, dirs, fnames in os.walk(proj_path):
        for f in fnames:
            fpath = os.path.join(root, f)
            rel = os.path.relpath(fpath, proj_path)
            files.append({"name": rel, "size": os.path.getsize(fpath)})
    return files


# --- Mobile project sync -------------------------------------------------
# The iOS app keeps an IndexedDB cache. These endpoints provide a compact
# manifest plus idempotent mutations so reconnecting never overwrites a newer
# local change silently.
def _sync_log_path(project_path: str) -> str:
    return os.path.join(project_path, ".mobile_sync_operations.json")


def _read_sync_log(project_path: str) -> set[str]:
    try:
        with open(_sync_log_path(project_path), "r", encoding="utf-8") as f:
            data = json.load(f)
        return set(data if isinstance(data, list) else [])
    except Exception:
        return set()


def _write_sync_log(project_path: str, operation_ids: set[str]) -> None:
    with open(_sync_log_path(project_path), "w", encoding="utf-8") as f:
        json.dump(list(operation_ids)[-5000:], f, ensure_ascii=False)


def _mobile_file_path(project_path: str, filename: str) -> str:
    clean = str(filename or "").replace("\\", "/").lstrip("/")
    if not clean or "\x00" in clean or any(part in ("", ".", "..") for part in clean.split("/")):
        raise HTTPException(400, "Invalid file path")
    destination = os.path.realpath(os.path.join(project_path, clean))
    if not destination.startswith(os.path.realpath(project_path) + os.sep):
        raise HTTPException(403, "Path traversal denied")
    return destination


@app.get("/api/projects/{name}/sync")
async def get_mobile_sync_manifest(name: str):
    project_path = os.path.join(_get_projects_dir(), name)
    if not os.path.isdir(project_path):
        raise HTTPException(404, "Project not found")
    files = []
    for root, _, filenames in os.walk(project_path):
        for filename in filenames:
            if filename == ".mobile_sync_operations.json":
                continue
            full_path = os.path.join(root, filename)
            relative = os.path.relpath(full_path, project_path).replace(os.sep, "/")
            files.append({
                "path": relative,
                "size": os.path.getsize(full_path),
                "mtime": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime(os.path.getmtime(full_path))),
            })
    return {"project": name, "files": sorted(files, key=lambda item: item["path"])}


@app.post("/api/projects/{name}/sync")
async def apply_mobile_sync_changes(name: str, request: Request):
    body = await request.json()
    changes = body.get("changes", []) if isinstance(body, dict) else []
    if not isinstance(changes, list):
        raise HTTPException(400, "changes must be a list")
    project_path = os.path.join(_get_projects_dir(), name)
    applied: list[str] = []

    # A first local sync may create the remote project before its file writes.
    create_change = next((item for item in changes if isinstance(item, dict) and item.get("type") == "create_project"), None)
    if not os.path.isdir(project_path):
        if not create_change:
            raise HTTPException(404, "Project not found")
        config = create_change.get("config") or {}
        os.makedirs(project_path, exist_ok=True)
        _create_default_knowledge_tree(project_path)
        _create_study_plan(
            project_path, config.get("exam_date", ""), config.get("exam_name", name),
            config.get("exam_type", ""), config.get("province", ""), config.get("mock_exam_count", 120),
            config.get("business_model") or {},
        )

    operation_ids = _read_sync_log(project_path)
    for change in changes:
        if not isinstance(change, dict):
            continue
        operation_id = str(change.get("id") or change.get("idempotency_key") or "")
        if operation_id and operation_id in operation_ids:
            applied.append(operation_id)
            continue
        kind = change.get("type")
        if kind == "create_project":
            pass
        elif kind == "write":
            path = _mobile_file_path(project_path, change.get("path", ""))
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "w", encoding="utf-8") as f:
                f.write(str(change.get("content", "")))
        elif kind == "delete":
            path = _mobile_file_path(project_path, change.get("path", ""))
            if os.path.isfile(path):
                os.remove(path)
        else:
            raise HTTPException(400, "Unsupported sync change")
        if operation_id:
            operation_ids.add(operation_id)
            applied.append(operation_id)
    _write_sync_log(project_path, operation_ids)
    return {"ok": True, "applied": applied}


@app.delete("/api/projects/{name}/files/")
async def clear_project_data(name: str):
    """Clear all practice data files but keep config files (备考计划.json, 能力画像.json, etc)."""
    proj_path = os.path.join(_get_projects_dir(), name)
    if not os.path.isdir(proj_path):
        raise HTTPException(404, "Project not found")
    keep = {'备考计划.json', '能力画像.json', '练习统计.json', '知识体系.json'}
    deleted = 0
    for root, dirs, fnames in os.walk(proj_path):
        # Skip top-level config files
        if root == proj_path:
            dirs[:] = [d for d in dirs if d in ('练习', '每日热点', '每日知识点')]
            continue
        for f in fnames:
            fpath = os.path.join(root, f)
            os.remove(fpath)
            deleted += 1
    # Clean up empty dirs
    for root, dirs, fnames in os.walk(proj_path, topdown=False):
        for d in dirs:
            dp = os.path.join(root, d)
            try:
                if not os.listdir(dp):
                    os.rmdir(dp)
            except:
                pass
    return {"message": f"已清除 {deleted} 个文件"}


def _count_cases(path: str) -> int:
    """Count test cases in a JSON/JSONL file."""
    if not os.path.exists(path):
        return 0
    try:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read().strip()
        if not content:
            return 0
        # Try JSON array first
        if content.startswith("["):
            data = json.loads(content)
            return len(data) if isinstance(data, list) else 1
        # Try JSON wrapper (single object with possible test_cases key)
        # But JSONL also starts with {, so we need to handle parse errors
        if content.startswith("{"):
            try:
                data = json.loads(content)
                if isinstance(data, list):
                    return len(data)
                if "test_cases" in data:
                    return len(data["test_cases"])
                return 1
            except json.JSONDecodeError:
                # Not valid single JSON — treat as JSONL
                pass
        # JSONL: one JSON object per line
        return sum(1 for line in content.split("\n") if line.strip())
    except Exception:
        return 0


def _save_project_cases(path: str, cases: list[dict]):
    """Save test cases in JSONL format (one JSON object per line)."""
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        for tc in cases:
            f.write(json.dumps(tc, ensure_ascii=False) + "\n")


def _load_project_cases(name: str) -> list[dict]:
    """Load test cases from a project directory. Supports JSON array, JSON wrapper, and JSONL formats."""
    cases_dir = os.path.join(_get_projects_dir(), name, "cases")
    cases_path = os.path.join(cases_dir, "test_cases.jsonl")
    if not os.path.exists(cases_path):
        cases_path = os.path.join(cases_dir, "test_cases.json")
    if os.path.exists(cases_path):
        try:
            with open(cases_path, encoding="utf-8") as f:
                content = f.read().strip()
            if not content:
                return []

            # JSON array: [{...}, {...}]
            if content.startswith("["):
                data = json.loads(content)
                if isinstance(data, list):
                    return data
                if isinstance(data, dict):
                    return data.get("test_cases", [])

            # JSON wrapper or JSONL
            if content.startswith("{"):
                try:
                    data = json.loads(content)
                except json.JSONDecodeError:
                    # Not a single JSON object → try JSONL
                    data = None

                if isinstance(data, dict) and "test_cases" in data:
                    return data["test_cases"]
                if isinstance(data, dict):
                    # Single wrapped object
                    return [data]
                if isinstance(data, list):
                    return data

                # JSONL: one JSON object per line
                cases = []
                for line in content.split("\n"):
                    line = line.strip()
                    if line:
                        try:
                            cases.append(json.loads(line))
                        except json.JSONDecodeError:
                            pass
                return cases

        except (IOError,):
            pass
    return []


def _list_spec_files(name: str) -> list[str]:
    """List spec files in a project."""
    spec_dir = os.path.join(_get_projects_dir(), name, "spec")
    if not os.path.exists(spec_dir):
        return []
    return [f for f in os.listdir(spec_dir) if os.path.isfile(os.path.join(spec_dir, f))]


# Static files mount (lucide.js, etc.) — no cache for mobile dev
_static_dir = os.path.join(os.path.dirname(__file__), "static")
app.mount("/static", StaticFiles(directory=_static_dir, html=True), name="static")

# Add cache-control to all static responses
from starlette.middleware.base import BaseHTTPMiddleware
class NoCacheMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request, call_next):
        response = await call_next(request)
        if request.url.path.startswith("/static/"):
            response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
            response.headers["Pragma"] = "no-cache"
            response.headers["Expires"] = "0"
        return response
app.add_middleware(NoCacheMiddleware)

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8765))
    uvicorn.run(app, host="0.0.0.0", port=port)
